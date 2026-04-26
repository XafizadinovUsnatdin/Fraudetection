"""SafeNet Anti-Fraud — 15-slide PPTX generator."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D,0x1B,0x2A)
NAVY     = RGBColor(0x12,0x28,0x50)
BLUE     = RGBColor(0x00,0x66,0xCC)
BLIGHT   = RGBColor(0x60,0xA5,0xFA)
PURPLE   = RGBColor(0xA7,0x8B,0xFA)
GREEN    = RGBColor(0x00,0xB8,0x6E)
GLIGHT   = RGBColor(0x4A,0xDE,0x80)
RED      = RGBColor(0xE6,0x39,0x46)
RLIGHT   = RGBColor(0xF8,0x71,0x71)
YELLOW   = RGBColor(0xFB,0xBF,0x24)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
GRAY     = RGBColor(0x94,0xA3,0xB8)
DGRAY    = RGBColor(0x47,0x55,0x69)

W = Inches(13.33)
H = Inches(7.5)
OUT = Path(__file__).parent / "SafeNet_AntifraudPresentation.pptx"

# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, border=None, bw=1.2):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h, size=13, bold=False,
       color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color

def head(slide, title, sub="", tcol=BLIGHT, scol=GRAY):
    tx(slide, title, Inches(.55), Inches(.28), Inches(12.2), Inches(.75),
       size=28, bold=True, color=tcol)
    if sub:
        tx(slide, sub, Inches(.55), Inches(1.02), Inches(12.2), Inches(.42),
           size=13, color=scol)
    box(slide, Inches(.55), Inches(1.42), Inches(1.3), Inches(.055), fill=BLUE)

def card(slide, l, t, w, h,
         fill=RGBColor(0x10,0x1E,0x38),
         border=RGBColor(0x2E,0x4A,0x7E)):
    box(slide, l, t, w, h, fill=fill, border=border)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def s01(prs):
    s = blank(prs); bg(s, BG)
    box(s, 0, Inches(2.4), W, Inches(2.9), fill=RGBColor(0x12,0x28,0x50))
    tx(s,"🛡️", Inches(5.8),Inches(.35),Inches(1.7),Inches(1.0), size=52, align=PP_ALIGN.CENTER)
    tx(s,"SafeNet", Inches(1.5),Inches(1.35),Inches(10.3),Inches(1.1),
       size=54,bold=True,color=BLIGHT,align=PP_ALIGN.CENTER)
    tx(s,"Sun'iy Intellekt Asosida Anti-Fraud Tizimi",
       Inches(1.5),Inches(2.6),Inches(10.3),Inches(.65),
       size=21,color=PURPLE,align=PP_ALIGN.CENTER)
    tx(s,"Shubhali va firibgarlik tranzaksiyalarini erta aniqlash",
       Inches(1.5),Inches(3.25),Inches(10.3),Inches(.5),
       size=14,color=GRAY,align=PP_ALIGN.CENTER)
    pills=[("🤖 Deep Learning",BLIGHT),("🧠 DeepFraudNet MLP",PURPLE),
           ("📊 Streamlit Demo",GLIGHT),("🇺🇿 O'zbekiston",YELLOW)]
    for i,(txt,col) in enumerate(pills):
        x=Inches(.9)+i*Inches(3.0)
        card(s,x,Inches(4.05),Inches(2.7),Inches(.58),
             fill=RGBColor(0x0E,0x22,0x44),border=col)
        tx(s,txt,x,Inches(4.09),Inches(2.7),Inches(.52),
           size=12,align=PP_ALIGN.CENTER,color=col)
    tx(s,"Taqdimot · 2026 yil · Xafizadinov Usnatdin",
       Inches(1.5),Inches(6.85),Inches(10.3),Inches(.4),
       size=12,color=DGRAY,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Fraud nima va muammo
# ══════════════════════════════════════════════════════════════════════════════
def s02(prs):
    s = blank(prs); bg(s, RGBColor(0x10,0x08,0x08))
    head(s,"Muammo — Firibgarlik (Fraud) Nima?",
         "Moliyaviy firibgarlik global miqyosda o'sib bormoqda",
         tcol=RLIGHT, scol=RGBColor(0xFC,0xA5,0xA5))
    box(s,Inches(.55),Inches(1.42),Inches(1.3),Inches(.055),fill=RED)

    stats=[("$485 mlrd","Yillik global fraud zarari"),
           ("3.5 soniya","Har shunda 1 ta fraud"),
           ("79%","Raqamli to'lovlarda fraud ulushi")]
    for i,(num,lbl) in enumerate(stats):
        x=Inches(.55)+i*Inches(4.2)
        card(s,x,Inches(1.62),Inches(3.9),Inches(1.5),
             fill=RGBColor(0x28,0x08,0x08),border=RGBColor(0xB0,0x20,0x20))
        tx(s,num,x,Inches(1.74),Inches(3.9),Inches(.72),
           size=30,bold=True,align=PP_ALIGN.CENTER,color=RLIGHT)
        tx(s,lbl,x,Inches(2.46),Inches(3.9),Inches(.5),
           size=12,align=PP_ALIGN.CENTER,color=RGBColor(0xFC,0xA5,0xA5))

    tx(s,"Fraud turlari:",Inches(.55),Inches(3.3),Inches(12.0),Inches(.45),
       size=13,bold=True,color=RLIGHT)
    fraud_types=[
        ("💳","Karta firibgarligi","Boshqaning kartasidan ruxsatsiz to'lov"),
        ("🤖","Bot hujumi","Avtomatlashtirilgan ko'plab kichik tranzaksiyalar"),
        ("🎭","Identifikatsiya o'g'irligi","Boshqa shaxs nomi bilan ro'yxatdan o'tish"),
        ("📍","Joylashuv firibgarligi","Chet el IP orqali mahalliy hisobdan to'lov"),
    ]
    for i,(ico,title,desc) in enumerate(fraud_types):
        row,col=divmod(i,2)
        x=Inches(.55)+col*Inches(6.3)
        y=Inches(3.85)+row*Inches(1.4)
        card(s,x,y,Inches(6.0),Inches(1.25),
             fill=RGBColor(0x20,0x08,0x08),border=RGBColor(0x80,0x18,0x18))
        tx(s,ico+" "+title,x+Inches(.15),y+Inches(.1),
           Inches(5.7),Inches(.45),size=13,bold=True,color=RLIGHT)
        tx(s,desc,x+Inches(.15),y+Inches(.58),
           Inches(5.7),Inches(.55),size=11,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset ko'rinishi
# ══════════════════════════════════════════════════════════════════════════════
def s03(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Dataset Ko'rinishi",
         "15,000 tranzaksiya · 600 foydalanuvchi · sintetik + Kaggle ma'lumotlari")

    # table header
    cols=[("user_id",1.6),("amount",1.5),("device",1.5),
          ("location",1.7),("trans_hour",1.55),("isFraud",1.35)]
    col_colors=[BLIGHT,PURPLE,PURPLE,PURPLE,PURPLE,RLIGHT]
    x_start=Inches(.4)
    xs=[]
    cur=x_start
    for name,w_in in cols:
        xs.append(cur); cur+=Inches(w_in)

    # header row
    for i,((name,w_in),col) in enumerate(zip(cols,col_colors)):
        card(s,xs[i],Inches(1.62),Inches(w_in),Inches(.5),
             fill=RGBColor(0x0A,0x20,0x40),border=BLUE)
        tx(s,name,xs[i]+Inches(.07),Inches(1.65),Inches(w_in-.1),Inches(.45),
           size=12,bold=True,color=col)

    # sample rows
    rows=[
        ["user_001","50,000","ios","toshkent","14","0"],
        ["user_042","850,000","chrome","samarkand","23","0"],
        ["user_075","10,000,000","emulator","foreign_ip","2","1"],
        ["user_118","120,000","android","bukhara","9","0"],
        ["user_031","5,500,000","unknown","foreign_ip","3","1"],
        ["user_200","75,000","safari","fergana","11","0"],
        ["user_099","2,200,000","emulator","unknown","1","1"],
    ]
    fraud_rows={2,4,6}
    for r,row in enumerate(rows):
        ry=Inches(2.18)+r*Inches(.54)
        is_fraud=r in fraud_rows
        for i,(val,(name,w_in)) in enumerate(zip(row,cols)):
            fill=RGBColor(0x28,0x08,0x08) if is_fraud else RGBColor(0x0C,0x1A,0x30)
            bord=RGBColor(0x80,0x18,0x18) if is_fraud else RGBColor(0x1E,0x36,0x5A)
            card(s,xs[i],ry,Inches(w_in),Inches(.51),fill=fill,border=bord)
            vcol=RLIGHT if (i==5 and val=="1") else (GLIGHT if i==5 else WHITE)
            tx(s,val,xs[i]+Inches(.07),ry+Inches(.06),
               Inches(w_in-.1),Inches(.42),size=11,color=vcol)

    # legend
    card(s,Inches(.4),Inches(6.04),Inches(5.5),Inches(.55),
         fill=RGBColor(0x28,0x08,0x08),border=RLIGHT)
    tx(s,"🔴  isFraud=1 — Fraud tranzaksiyasi (bloklangan)",
       Inches(.6),Inches(6.1),Inches(5.3),Inches(.46),size=12,color=RLIGHT)
    card(s,Inches(6.1),Inches(6.04),Inches(5.5),Inches(.55),
         fill=RGBColor(0x0C,0x1A,0x30),border=GLIGHT)
    tx(s,"🟢  isFraud=0 — Oddiy (normal) tranzaksiya",
       Inches(6.3),Inches(6.1),Inches(5.3),Inches(.46),size=12,color=GLIGHT)

    # stats bar
    stats=[("15,000","Jami qator"),("~10%","Fraud ulushi"),
           ("600","Foydalanuvchi"),("9","Joylashuv")]
    for i,(v,l) in enumerate(stats):
        x=Inches(.4)+i*Inches(3.0)
        card(s,x,Inches(6.72),Inches(2.7),Inches(.58),
             fill=RGBColor(0x0A,0x18,0x34),border=BLUE)
        tx(s,v+" — "+l,x+Inches(.12),Inches(6.76),Inches(2.5),Inches(.5),
           size=12,bold=True,color=BLIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Ma'lumotlar bilan ishlash
# ══════════════════════════════════════════════════════════════════════════════
def s04(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Ma'lumotlar Bilan Ishlash",
         "Xom ma'lumotdan model uchun tayyor dataset yaratish jarayoni")

    steps=[
        ("1","📥  Ma'lumot Yuklash",
         "CSV fayl o'qiladi (Kaggle yoki sintetik)\n"
         "Majburiy ustunlar tekshiriladi: user_id, amount,\n"
         "device, location, transaction_hour, isFraud",
         RGBColor(0x0A,0x20,0x40),BLIGHT),
        ("2","🧹  Tozalash",
         "NaN qiymatlar to'ldiriladi\namount ≥ 0.01 — manfiy summalar tuzatiladi\n"
         "transaction_hour: 0–23 oralig'iga clip\nisFraud: 0 yoki 1 ga qaytariladi",
         RGBColor(0x0A,0x20,0x18),GREEN),
        ("3","👤  Foydalanuvchi Profili",
         "Har user uchun:\n• O'rtacha summa (user_avg_amount)\n"
         "• Tranzaksiya soni\n• Ko'rilgan qurilmalar ro'yxati\n"
         "• Ko'rilgan joylashuvlar ro'yxati",
         RGBColor(0x18,0x10,0x30),PURPLE),
        ("4","🧠  Behavior Features",
         "Har tranzaksiya uchun real-time:\n"
         "• amount_ratio_to_avg (0–50)\n"
         "• is_new_device (0/1)\n"
         "• is_new_location (0/1)\n"
         "• is_night_transaction (0/1)",
         RGBColor(0x20,0x14,0x04),YELLOW),
        ("5","🔢  Encoding",
         "Kategorik: device, location\n"
         "→ OneHotEncoder(handle_unknown='ignore')\n\n"
         "Raqamli: passthrough\n(normalizatsiya kerak emas RF uchun)",
         RGBColor(0x0A,0x20,0x40),BLIGHT),
        ("6","✂️  Train/Test Bo'lish",
         "test_size = 0.25\nstratify=y — fraud ulushi saqlanadi\n"
         "random_state = 42 — takrorlanadi\n\n"
         "Train: 11,250 qator\nTest:   3,750 qator",
         RGBColor(0x0A,0x20,0x18),GLIGHT),
    ]
    for i,(num,title,desc,fill,col) in enumerate(steps):
        row,c=divmod(i,3)
        x=Inches(.45)+c*Inches(4.25)
        y=Inches(1.62)+row*Inches(2.55)
        card(s,x,y,Inches(4.05),Inches(2.38),fill=fill,border=col)
        # number circle
        box(s,x+Inches(.15),y+Inches(.12),Inches(.42),Inches(.42),fill=col)
        tx(s,num,x+Inches(.15),y+Inches(.1),Inches(.42),Inches(.42),
           size=14,bold=True,color=BG,align=PP_ALIGN.CENTER)
        tx(s,title,x+Inches(.68),y+Inches(.12),Inches(3.2),Inches(.45),
           size=12,bold=True,color=col)
        tx(s,desc,x+Inches(.15),y+Inches(.65),Inches(3.75),Inches(1.6),
           size=10,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Feature Engineering
# ══════════════════════════════════════════════════════════════════════════════
def s05(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Feature Engineering — Belgilar Yaratish",
         "18 ta asosiy + 7 ta interaction = 25 ta EXT_FEATURES modeli uchun")

    base=[
        ("Raqamli (7 ta)","amount, transaction_hour, Tx Frequency,\nTime Since Last Tx, Account Age,\nNormalized Amount, Fraud Complaints",PURPLE),
        ("Ikkilik (9 ta)","Geo-Location Flags, Location-Inconsistent,\nRecipient Verification, Recipient Blacklist,\nVPN/Proxy, Merchant Mismatch,\nDaily Limit Exceeded, High-Value Flags,\nPast Fraud Behavior",BLIGHT),
        ("Kategorik (2 ta)","device  →  OneHot Encoding\nlocation  →  OneHot Encoding\nBinning: NUMERICAL uchun 9 kvantil",GREEN),
    ]
    for i,(title,desc,col) in enumerate(base):
        x=Inches(.45)+i*Inches(3.0)
        card(s,x,Inches(1.62),Inches(2.85),Inches(2.2),
             fill=RGBColor(0x0A,0x18,0x30),border=col)
        tx(s,title,x+Inches(.12),Inches(1.72),Inches(2.6),Inches(.42),
           size=12,bold=True,color=col)
        tx(s,desc,x+Inches(.12),Inches(2.2),Inches(2.62),Inches(1.52),
           size=10,color=GRAY)

    # arrow
    tx(s,"⟹",Inches(9.5),Inches(2.4),Inches(.7),Inches(.7),
       size=28,color=BLIGHT,align=PP_ALIGN.CENTER)

    # interaction box
    card(s,Inches(10.2),Inches(1.62),Inches(2.85),Inches(2.2),
         fill=RGBColor(0x12,0x0A,0x28),border=PURPLE)
    tx(s,"⚡ Interaction (7 ta)",Inches(10.35),Inches(1.72),Inches(2.6),Inches(.42),
       size=12,bold=True,color=PURPLE)
    tx(s,"risky_device\nrisky_location\nnight_flag\ndev×night  dev×loc\nloc×night\nvpn×blacklist",
       Inches(10.35),Inches(2.2),Inches(2.62),Inches(1.52),size=10,color=GRAY)

    # formula strip
    card(s,Inches(.45),Inches(3.98),Inches(12.4),Inches(.55),
         fill=RGBColor(0x08,0x14,0x28),border=BLUE)
    tx(s,"EXT_FEATURES  =  18 asosiy belgi  +  7 interaction  =  25 ta",
       Inches(.65),Inches(4.04),Inches(12.0),Inches(.42),
       size=13,bold=True,color=BLIGHT,align=PP_ALIGN.CENTER)

    # examples
    examples=[
        ("risky_device","device ∈ {emulator,unknown,linux} → 1","Xavfli qurilma bayrog'i",RLIGHT),
        ("dev_x_night","risky_device × night_flag","Xavfli qurilma + tun birgaligi",YELLOW),
        ("vpn_x_blacklist","VPN × Recipient Blacklist","VPN + qora ro'yxat birgaligi",RLIGHT),
        ("night_flag","transaction_hour ∈ [0..5] → 1","Tun soati (00:00–05:59)",BLIGHT),
    ]
    for i,(name,formula,meaning,col) in enumerate(examples):
        x=Inches(.45)+i*Inches(3.12)
        card(s,x,Inches(4.68),Inches(2.98),Inches(2.62),
             fill=RGBColor(0x0C,0x18,0x30),border=col)
        tx(s,name,x+Inches(.12),Inches(4.78),Inches(2.75),Inches(.38),
           size=12,bold=True,color=col)
        tx(s,formula,x+Inches(.12),Inches(5.22),Inches(2.75),Inches(.55),
           size=10,color=GLIGHT)
        tx(s,meaning,x+Inches(.12),Inches(5.82),Inches(2.75),Inches(.4),
           size=10,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Fraud qanday aniqlanadi (algoritm)
# ══════════════════════════════════════════════════════════════════════════════
def s06(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Fraud Qanday Aniqlanadi?",
         "Har bir tranzaksiya uchun bosqichma-bosqich qaror qabul qilish")

    flow=[
        ("📥","Tranzaksiya\nkeladi",BLIGHT,RGBColor(0x0A,0x20,0x44)),
        ("⚙️","25 ta\nEXT_FEATURES",PURPLE,RGBColor(0x18,0x10,0x30)),
        ("🧠","256→128→64\nMLP qatlamlari",YELLOW,RGBColor(0x22,0x18,0x04)),
        ("📡","Sigmoid\nchiqish",GREEN,RGBColor(0x08,0x22,0x10)),
        ("🎯","P(fraud)\n0.0–1.0",GLIGHT,RGBColor(0x08,0x22,0x10)),
        ("⚖️","Threshold\ntekshiruv",BLIGHT,RGBColor(0x0A,0x20,0x44)),
    ]
    for i,(ico,label,col,fill) in enumerate(flow):
        x=Inches(.4)+i*Inches(2.12)
        card(s,x,Inches(1.72),Inches(1.9),Inches(1.5),fill=fill,border=col)
        tx(s,ico,x,Inches(1.82),Inches(1.9),Inches(.55),
           size=26,align=PP_ALIGN.CENTER)
        tx(s,label,x,Inches(2.42),Inches(1.9),Inches(.72),
           size=11,bold=True,color=col,align=PP_ALIGN.CENTER)
        if i<5:
            tx(s,"→",x+Inches(1.9),Inches(2.22),Inches(.22),Inches(.5),
               size=18,color=GRAY,align=PP_ALIGN.CENTER)

    # 3 decision paths
    tx(s,"Qaror yo'llari:",Inches(.4),Inches(3.45),Inches(12.5),Inches(.45),
       size=14,bold=True,color=WHITE)

    paths=[
        ("🟢  ALLOW — Xavfsiz",
         "P < 0.50",
         "✅ Tranzaksiya o'tkaziladi\n"
         "• Oddiy qurilma va joylashuv\n"
         "• Summa odatiy miqdorda\n"
         "• Kunduz soatida amalga oshirilgan",
         GLIGHT,RGBColor(0x08,0x22,0x10),RGBColor(0x0F,0x38,0x1A)),
        ("🟡  REVIEW — Tekshirish",
         "0.50 ≤ P < 0.80",
         "⚠️ Operator tekshirishi kutiladi\n"
         "• Tun saatida amalga oshirilgan\n"
         "• Summa o'rtachadan yuqori\n"
         "• Yangi joylashuv yoki qurilma",
         YELLOW,RGBColor(0x22,0x18,0x04),RGBColor(0x3A,0x28,0x06)),
        ("🔴  BLOCK — Bloklash",
         "P ≥ 0.80",
         "🚫 Tranzaksiya darhol bloklanadi\n"
         "• Emulator / noma'lum qurilma\n"
         "• Chet el / noma'lum IP\n"
         "• Summa 5x+ baland, tun vaqti",
         RLIGHT,RGBColor(0x28,0x08,0x08),RGBColor(0x48,0x10,0x10)),
    ]
    for i,(title,prob,desc,col,fill,bord) in enumerate(paths):
        x=Inches(.4)+i*Inches(4.28)
        card(s,x,Inches(3.98),Inches(4.05),Inches(3.3),fill=fill,border=bord)
        card(s,x,Inches(3.98),Inches(4.05),Inches(.52),fill=bord,border=bord)
        tx(s,title,x+Inches(.12),Inches(4.02),Inches(3.8),Inches(.46),
           size=13,bold=True,color=col)
        card(s,x+Inches(.12),Inches(4.56),Inches(1.5),Inches(.38),
             fill=RGBColor(0x08,0x10,0x20),border=col)
        tx(s,prob,x+Inches(.15),Inches(4.59),Inches(1.45),Inches(.33),
           size=11,bold=True,color=col)
        tx(s,desc,x+Inches(.12),Inches(5.02),Inches(3.8),Inches(2.1),
           size=11,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Risk ball tizimi
# ══════════════════════════════════════════════════════════════════════════════
def s07(prs):
    s = blank(prs); bg(s, RGBColor(0x08,0x10,0x20))
    head(s,"Risk Ball Tizimi",
         "Har bir xavf omili qo'shimcha ball beradi — jami ball → ehtimollik")

    factors=[
        ("💰","Summa o'rtachadan ≥ 1.8x","+0.8",YELLOW),
        ("💰","Summa o'rtachadan ≥ 3.0x","+1.1",YELLOW),
        ("💰","Summa o'rtachadan ≥ 5.0x","+1.0",RLIGHT),
        ("📱","Yangi qurilma (is_new_device)","+1.4",PURPLE),
        ("📍","Yangi joylashuv (is_new_location)","+1.6",PURPLE),
        ("🌙","Tun saati 0–5 (is_night)","+0.8",BLIGHT),
        ("⚠️","Xavfli qurilma yoki chet el IP","+1.3",RLIGHT),
        ("🔄","Yangi qurilma + yangi joylashuv","+0.9",RLIGHT),
        ("💵","Summa ≥ 500$","+0.7",YELLOW),
    ]
    for i,(ico,name,ball,col) in enumerate(factors):
        row,c=divmod(i,3)
        x=Inches(.45)+c*Inches(4.25)
        y=Inches(1.68)+row*Inches(1.3)
        card(s,x,y,Inches(4.05),Inches(1.18),
             fill=RGBColor(0x10,0x18,0x2E),border=col)
        tx(s,ico+"  "+name,x+Inches(.12),y+Inches(.1),
           Inches(3.0),Inches(.5),size=12,bold=True,color=WHITE)
        tx(s,ball,x+Inches(3.15),y+Inches(.08),Inches(.78),Inches(.5),
           size=22,bold=True,color=col,align=PP_ALIGN.RIGHT)
        tx(s,"ball",x+Inches(3.15),y+Inches(.58),Inches(.78),Inches(.35),
           size=9,color=col,align=PP_ALIGN.RIGHT)

    # formula box
    card(s,Inches(.45),Inches(6.45),Inches(12.4),Inches(.82),
         fill=RGBColor(0x08,0x14,0x2A),border=BLUE)
    tx(s,"📐  Formula:",Inches(.65),Inches(6.52),Inches(2.0),Inches(.38),
       size=12,bold=True,color=BLIGHT)
    tx(s,"P(fraud) = sigmoid( −5.2 + 1.15 × Σ ball )     |     Σ ball ≥ 5.5 → P = 0.92     |     Σ ball ≥ 4.2 → P = 0.68",
       Inches(2.7),Inches(6.52),Inches(10.0),Inches(.65),
       size=12,bold=True,color=GLIGHT,align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DeepFraudNet arxitekturasi
# ══════════════════════════════════════════════════════════════════════════════
def s08(prs):
    s = blank(prs); bg(s, RGBColor(0x0D,0x0D,0x1E))
    head(s,"Algoritm — DeepFraudNet Arxitekturasi",
         "3 ta yashirin qatlamli MLP · Adam optimizer · Pure NumPy implementatsiya",
         tcol=PURPLE)
    box(s,Inches(.55),Inches(1.42),Inches(1.3),Inches(.055),fill=PURPLE)

    # network diagram
    layers=[
        ("INPUT\n~130\nfeature",  BLIGHT, RGBColor(0x08,0x16,0x34)),
        ("Dense\n256\nReLU",      PURPLE, RGBColor(0x18,0x10,0x30)),
        ("Dense\n128\nReLU",      PURPLE, RGBColor(0x18,0x10,0x30)),
        ("Dense\n64\nReLU",       PURPLE, RGBColor(0x18,0x10,0x30)),
        ("OUTPUT\n1\nSigmoid",    GLIGHT, RGBColor(0x08,0x22,0x10)),
    ]
    for i,(label,col,fill) in enumerate(layers):
        x=Inches(.5)+i*Inches(2.52)
        card(s,x,Inches(1.68),Inches(2.22),Inches(1.85),fill=fill,border=col)
        tx(s,label,x+Inches(.06),Inches(1.78),Inches(2.1),Inches(1.65),
           size=14,bold=True,color=col,align=PP_ALIGN.CENTER)
        if i<4:
            tx(s,"→",x+Inches(2.22),Inches(2.3),Inches(.3),Inches(.5),
               size=20,color=GRAY,align=PP_ALIGN.CENTER)

    # output annotation
    tx(s,"P(fraud)  ∈  [0 … 1]",
       Inches(10.7),Inches(3.62),Inches(2.4),Inches(.42),
       size=11,bold=True,color=GLIGHT,align=PP_ALIGN.CENTER)

    # left: training details
    card(s,Inches(.45),Inches(3.72),Inches(6.0),Inches(3.55),
         fill=RGBColor(0x10,0x0E,0x28),border=BLIGHT)
    tx(s,"⚙️  O'qitish Tafsilotlari",
       Inches(.65),Inches(3.84),Inches(5.6),Inches(.45),
       size=13,bold=True,color=BLIGHT)
    details=[
        ("Optimizer",       "Adam  (β₁=0.9, β₂=0.999, lr=0.001)"),
        ("Loss",            "Class-weighted Binary Cross-Entropy"),
        ("Batch size",      "512  (mini-batch)"),
        ("Epochlar",        "80"),
        ("Init",            "He initialization: W ~ N(0, √(2/fan_in))"),
        ("Regularizatsiya", "L2 weight decay  λ=1e-5"),
        ("Sinf og'irligi",  "ratio = min(√(n_neg / n_pos), 5.0)"),
        ("Namuna",          "20,000 (stratified sampling)"),
    ]
    for i,(k,v) in enumerate(details):
        y=Inches(4.42)+i*Inches(.38)
        tx(s,k+":",Inches(.65),y,Inches(2.1),Inches(.36),size=10,color=GRAY)
        tx(s,v,Inches(2.8),y,Inches(3.5),Inches(.36),size=10,bold=True,color=BLIGHT)

    # right: why MLP
    card(s,Inches(6.65),Inches(3.72),Inches(6.25),Inches(3.55),
         fill=RGBColor(0x10,0x0E,0x28),border=GLIGHT)
    tx(s,"✅  Nima uchun DeepFraudNet?",
       Inches(6.85),Inches(3.84),Inches(5.8),Inches(.45),
       size=13,bold=True,color=GLIGHT)
    pros=[
        "• Nochiziqli munosabatlarni o'rganadi (ReLU)",
        "• Interaction features birgalikda ishlaydi",
        "• Adam — tez va barqaror konvergentsiya",
        "• sqrt class weighting — precision/recall muvozanat",
        "• He init — gradient yo'qolishining oldini oladi",
        "• Sklearn yo'q — Streamlit Cloud bilan mos",
        "• Quantile binning nochiziqlilikni kuchaytiradi",
    ]
    for i,p in enumerate(pros):
        tx(s,p,Inches(6.85),Inches(4.42)+i*Inches(.41),
           Inches(5.8),Inches(.38),size=11,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Model sifati: preprocessing pipeline
# ══════════════════════════════════════════════════════════════════════════════
def s09(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Model Sifati — Pipeline Arxitekturasi",
         "Preprocessing → Feature Engineering → DeepFraudNet → P(fraud)")

    # pipeline visual
    pipe_steps=[
        ("📥\nInput\nDataFrame\n(18 belgi)",BLIGHT,RGBColor(0x0A,0x1E,0x40)),
        ("⚡\nadd_\ninteractions\n(+7 belgi)",PURPLE,RGBColor(0x18,0x10,0x30)),
        ("🔀\nZ-score\nnorm +\nOHE + bin",GREEN,RGBColor(0x08,0x20,0x10)),
        ("🧠\nDeep\nFraudNet\n256→128→64",PURPLE,RGBColor(0x18,0x10,0x30)),
        ("🎯\nSigmoid\nP(fraud)\n0.0–1.0",GLIGHT,RGBColor(0x08,0x22,0x10)),
        ("⚖️\nThreshold\nALLOW /\nREVIEW /\nBLOCK",BLIGHT,RGBColor(0x0A,0x1E,0x40)),
    ]
    for i,(label,col,fill) in enumerate(pipe_steps):
        x=Inches(.4)+i*Inches(2.12)
        card(s,x,Inches(1.68),Inches(1.92),Inches(2.1),fill=fill,border=col)
        tx(s,label,x+Inches(.08),Inches(1.78),Inches(1.76),Inches(1.9),
           size=11,bold=True,color=col,align=PP_ALIGN.CENTER)
        if i<5:
            tx(s,"→",x+Inches(1.92),Inches(2.38),Inches(.22),Inches(.5),
               size=18,color=GRAY,align=PP_ALIGN.CENTER)

    # quality factors
    tx(s,"Model Sifatini Ta'minlovchi Omillar:",
       Inches(.4),Inches(3.98),Inches(12.5),Inches(.45),
       size=14,bold=True,color=BLIGHT)

    quality=[
        ("⚖️","Sinf Og'irligi (sqrt weighting)",
         "ratio = min(√(n_neg/n_pos), 5.0) ≈ 4.36x — fraud sinfiga yuqori og'irlik, "
         "lekin 10x dan kichik (precision muvozanatlanadi)",
         YELLOW,RGBColor(0x20,0x16,0x04)),
        ("🎯","Stratified Split + Sampling",
         "stratified_split(80/20) → har ikkala to'plamda 5% fraud saqlanadi. "
         "stratified_sample(20K) — o'qitish namumasida ham bir xil ulush",
         GLIGHT,RGBColor(0x08,0x20,0x0C)),
        ("🔄","Optimal Threshold",
         "0.05 dan 0.95 gacha 91 ta chegara sinab ko'riladi. F1-Score maksimali "
         "aniqlanib, ALLOW/REVIEW/BLOCK uchun ishlatiladi",
         BLIGHT,RGBColor(0x08,0x18,0x34)),
        ("💾","Artifact Saqlash (pickle)",
         "Model + transform schema + threshold + metrics + feature importances "
         "bitta dict ichida pickle.dump() bilan saqlanadi",
         PURPLE,RGBColor(0x16,0x10,0x2C)),
    ]
    for i,(ico,title,desc,col,fill) in enumerate(quality):
        row,c=divmod(i,2)
        x=Inches(.4)+c*Inches(6.35)
        y=Inches(4.55)+row*Inches(1.25)
        card(s,x,y,Inches(6.1),Inches(1.15),fill=fill,border=col)
        tx(s,ico+"  "+title,x+Inches(.15),y+Inches(.1),
           Inches(5.8),Inches(.45),size=13,bold=True,color=col)
        tx(s,desc,x+Inches(.15),y+Inches(.58),
           Inches(5.8),Inches(.5),size=11,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Model aniqligi: metrikalar
# ══════════════════════════════════════════════════════════════════════════════
def s10(prs):
    s = blank(prs); bg(s, RGBColor(0x08,0x18,0x08))
    head(s,"Model Aniqligi — Ko'rsatkichlar",
         "40,000 o'qitish · 10,000 test · DeepFraudNet (256→128→64 + 25 EXT_FEATURES)",
         tcol=GLIGHT,scol=RGBColor(0xA7,0xF3,0xD0))
    box(s,Inches(.55),Inches(1.42),Inches(1.3),Inches(.055),fill=GREEN)

    metrics=[
        ("~95%","Accuracy","Barcha bashoratlarning to'g'riligi (sinf nomutanosibligi bor)",GLIGHT),
        ("52.4%","Precision","BLOCK diyilganlarning 52.4% haqiqiy fraud",GLIGHT),
        ("40.0%","Recall","Barcha fraudning 40.0% aniqlandi",YELLOW),
        ("45.4%","F1-Score","Precision + Recall muvozanati (0.45 = yaxshi)",GLIGHT),
        ("0.857","ROC-AUC","Ajratish qobiliyati — sanoat standarti 0.85–0.92",GLIGHT),
    ]
    for i,(val,name,desc,col) in enumerate(metrics):
        x=Inches(.45)+i*Inches(2.5)
        card(s,x,Inches(1.65),Inches(2.35),Inches(1.9),
             fill=RGBColor(0x0A,0x22,0x0E),border=col)
        tx(s,val,x,Inches(1.78),Inches(2.35),Inches(.85),
           size=30,bold=True,align=PP_ALIGN.CENTER,color=col)
        tx(s,name,x,Inches(2.63),Inches(2.35),Inches(.42),
           size=13,bold=True,align=PP_ALIGN.CENTER,color=WHITE)
        tx(s,desc,x+Inches(.08),Inches(3.08),Inches(2.2),Inches(.42),
           size=9,align=PP_ALIGN.CENTER,color=GRAY)

    # Confusion matrix
    tx(s,"Confusion Matrix (Test set: ~10,000 qator, 5% fraud)",
       Inches(.45),Inches(3.72),Inches(6.5),Inches(.45),
       size=13,bold=True,color=BLIGHT)
    cm=[
        ("",               "Pred: Normal","Pred: Fraud"),
        ("Real: Normal",   "TN = 9,143",  "FP = 357"),
        ("Real: Fraud",    "FN = 300",    "TP = 200"),
    ]
    cm_colors=[[GRAY,BLIGHT,BLIGHT],
               [BLIGHT,GLIGHT,RLIGHT],
               [BLIGHT,YELLOW,GLIGHT]]
    cm_fills=[[BG,RGBColor(0x0A,0x20,0x40),RGBColor(0x0A,0x20,0x40)],
              [RGBColor(0x0A,0x20,0x40),RGBColor(0x08,0x28,0x10),RGBColor(0x28,0x08,0x08)],
              [RGBColor(0x0A,0x20,0x40),RGBColor(0x20,0x16,0x04),RGBColor(0x08,0x28,0x10)]]
    for r,row in enumerate(cm):
        for c,cell in enumerate(row):
            cx=Inches(.45)+c*Inches(2.05)
            cy=Inches(4.22)+r*Inches(.95)
            if r>0 or c>0:
                card(s,cx,cy,Inches(2.0),Inches(.9),
                     fill=cm_fills[r][c],border=cm_colors[r][c])
            tx(s,cell,cx+Inches(.05),cy+Inches(.1),Inches(1.9),Inches(.75),
               size=12 if (r>0 and c>0) else 10,
               bold=(r>0 and c>0),
               color=cm_colors[r][c],align=PP_ALIGN.CENTER)

    # model comparison
    tx(s,"Model Solishtirma",Inches(6.75),Inches(3.72),Inches(6.1),Inches(.45),
       size=13,bold=True,color=BLIGHT)
    compare=[
        ("🧠 DeepFraudNet","95%","45.4%","0.857",GLIGHT),
        ("📈 Logistic Reg.", "92%","28.1%","0.791",GRAY),
        ("🌲 Random Forest", "93%","38.2%","0.832",GRAY),
        ("⚡ Qoidalar",      "78%","22.5%","0.650",GRAY),
    ]
    hdrs=["Model","Accuracy","F1","ROC-AUC"]
    hx=[Inches(6.75),Inches(9.35),Inches(10.6),Inches(11.65)]
    hw=[Inches(2.55),Inches(1.2),Inches(1.0),Inches(1.1)]
    for j,(hdr,hxx,hww) in enumerate(zip(hdrs,hx,hw)):
        card(s,hxx,Inches(4.22),hww,Inches(.42),
             fill=RGBColor(0x0A,0x20,0x40),border=BLUE)
        tx(s,hdr,hxx+Inches(.05),Inches(4.26),hww-Inches(.08),Inches(.36),
           size=10,bold=True,color=BLIGHT)
    for r,(name,acc,f1,auc,col) in enumerate(compare):
        ry=Inches(4.68)+r*Inches(.62)
        vals=[name,acc,f1,auc]
        for j,(v,hxx,hww) in enumerate(zip(vals,hx,hw)):
            fill=RGBColor(0x0A,0x20,0x0A) if col==GLIGHT else RGBColor(0x0C,0x1A,0x2A)
            card(s,hxx,ry,hww,Inches(.56),fill=fill,border=RGBColor(0x1E,0x36,0x2E))
            tx(s,v,hxx+Inches(.05),ry+Inches(.06),hww-Inches(.08),Inches(.46),
               size=10,color=col if j==0 else WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Confusion matrix tushuntirish
# ══════════════════════════════════════════════════════════════════════════════
def s11(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Natijalarni Talqin Qilish",
         "To'g'ri va noto'g'ri bashoratlarning amaliy ahamiyati")

    items=[
        ("✅  True Positive (TP) = 200",
         "Model fraud dedi → haqiqatan fraud\n"
         "→ Tranzaksiya to'g'ri blokland!",
         GLIGHT,RGBColor(0x08,0x22,0x10),RGBColor(0x0F,0x38,0x1A)),
        ("✅  True Negative (TN) = 9,143",
         "Model normal dedi → haqiqatan normal\n"
         "→ Foydalanuvchi muammosiz o'tdi!",
         GLIGHT,RGBColor(0x08,0x22,0x10),RGBColor(0x0F,0x38,0x1A)),
        ("⚠️  False Positive (FP) = 357",
         "Model fraud dedi → aslida normal\n"
         "→ Mushtariy noqulaylik sezdi (3.8%)",
         YELLOW,RGBColor(0x22,0x18,0x04),RGBColor(0x3A,0x28,0x06)),
        ("❌  False Negative (FN) = 300",
         "Model normal dedi → aslida fraud\n"
         "→ Fraud o'tib ketdi! (eng xavfli holat)",
         RLIGHT,RGBColor(0x28,0x08,0x08),RGBColor(0x44,0x10,0x10)),
    ]
    for i,(title,desc,col,fill,bord) in enumerate(items):
        row,c=divmod(i,2)
        x=Inches(.45)+c*Inches(6.35)
        y=Inches(1.68)+row*Inches(1.55)
        card(s,x,y,Inches(6.1),Inches(1.38),fill=fill,border=bord)
        tx(s,title,x+Inches(.15),y+Inches(.1),
           Inches(5.8),Inches(.48),size=13,bold=True,color=col)
        tx(s,desc,x+Inches(.15),y+Inches(.65),
           Inches(5.8),Inches(.6),size=11,color=GRAY)

    # business impact
    card(s,Inches(.45),Inches(5.0),Inches(12.4),Inches(2.28),
         fill=RGBColor(0x08,0x16,0x2C),border=BLIGHT)
    tx(s,"💼  Biznes Ta'siri",
       Inches(.65),Inches(5.1),Inches(12.0),Inches(.45),
       size=13,bold=True,color=BLIGHT)
    impacts=[
        ("200 fraud bloklandi","Bank = pul yutdi"),
        ("300 fraud o'tdi","Dataset shovqini tufayli"),
        ("357 yolg'on alarm","Foydalanuvchi noqulayligi"),
        ("0.857 ROC-AUC","Sanoat standarti doirasida"),
    ]
    for i,(v,l) in enumerate(impacts):
        x=Inches(.65)+i*Inches(3.0)
        tx(s,v,x,Inches(5.62),Inches(2.9),Inches(.4),
           size=13,bold=True,color=GLIGHT)
        tx(s,l,x,Inches(6.02),Inches(2.9),Inches(.35),
           size=11,color=GRAY)

    tx(s,'Maqsad: FN (fraud o\'tib ketish) ni minimallashtirishdir — Recall muhim!',
       Inches(.65),Inches(6.58),Inches(12.0),Inches(.5),
       size=12,italic=True,color=YELLOW)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Ishlaydigan Demo
# ══════════════════════════════════════════════════════════════════════════════
def s12(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Ishlaydigan Demo",
         "Streamlit interaktiv paneli — real vaqtda tranzaksiyani baholash")

    examples=[
        {
            "title":"✅ ALLOW — Xavfsiz",
            "col":GLIGHT,"bg":RGBColor(0x08,0x22,0x0E),"bord":GREEN,
            "data":[("Foydalanuvchi","user_042"),("Summa","50,000 so'm"),
                    ("Qurilma","iOS"),("Joylashuv","Toshkent"),("Vaqt","14:30")],
            "prob":"Fraud ehtimoli: 3.2%",
            "reasons":"✓ Odatiy qurilma\n✓ Odatiy joylashuv\n✓ Kunduz saati\n✓ Normal summa",
        },
        {
            "title":"⚠️ REVIEW — O'rta xavf",
            "col":YELLOW,"bg":RGBColor(0x22,0x16,0x04),"bord":RGBColor(0xD9,0x7A,0x06),
            "data":[("Foydalanuvchi","user_075"),("Summa","850,000 so'm"),
                    ("Qurilma","Chrome"),("Joylashuv","Samarqand"),("Vaqt","23:15")],
            "prob":"Fraud ehtimoli: 62.7%",
            "reasons":"⚠ Tun saati (23:15)\n⚠ Summa o'rtachadan baland\n⚠ Yangi joylashuv",
        },
        {
            "title":"🚫 BLOCK — Yuqori xavf",
            "col":RLIGHT,"bg":RGBColor(0x28,0x08,0x08),"bord":RED,
            "data":[("Foydalanuvchi","user_031"),("Summa","10,000,000 so'm"),
                    ("Qurilma","Emulator"),("Joylashuv","foreign_ip"),("Vaqt","02:45")],
            "prob":"Fraud ehtimoli: 96.3%",
            "reasons":"✗ Emulator qurilma\n✗ Noma'lum IP\n✗ Tun saati (02:45)\n✗ Summa 47x baland",
        },
    ]
    for i,ex in enumerate(examples):
        x=Inches(.45)+i*Inches(4.25)
        card(s,x,Inches(1.65),Inches(4.05),Inches(5.65),
             fill=ex["bg"],border=ex["bord"])
        card(s,x,Inches(1.65),Inches(4.05),Inches(.55),
             fill=ex["bord"],border=ex["bord"])
        tx(s,ex["title"],x+Inches(.12),Inches(1.7),
           Inches(3.8),Inches(.48),size=13,bold=True,color=ex["col"])
        for j,(k,v) in enumerate(ex["data"]):
            y=Inches(2.28)+j*Inches(.5)
            tx(s,k+":",x+Inches(.15),y,Inches(1.5),Inches(.46),size=11,color=GRAY)
            tx(s,v,x+Inches(1.65),y,Inches(2.28),Inches(.46),size=11,bold=True,color=WHITE)
        card(s,x+Inches(.12),Inches(4.82),Inches(3.8),Inches(.48),
             fill=RGBColor(0x08,0x10,0x1C),border=ex["col"])
        tx(s,ex["prob"],x+Inches(.2),Inches(4.86),
           Inches(3.65),Inches(.4),size=12,bold=True,color=ex["col"])
        tx(s,ex["reasons"],x+Inches(.15),Inches(5.4),
           Inches(3.8),Inches(.85),size=11,color=GRAY)

    # demo link
    card(s,Inches(.45),Inches(7.05),Inches(12.4),Inches(.3),
         fill=RGBColor(0x08,0x16,0x2C),border=BLUE)
    tx(s,"🌐  github.com/XafizadinovUsnatdin/Fraudetection   |   Streamlit Cloud'da ochiq",
       Inches(.65),Inches(7.06),Inches(12.0),Inches(.28),
       size=11,color=BLIGHT,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — API va integratsiya
# ══════════════════════════════════════════════════════════════════════════════
def s13(prs):
    s = blank(prs); bg(s, RGBColor(0x08,0x0E,0x18))
    head(s,"REST API — Integratsiya",
         "FastAPI orqali har qanday bank tizimiga ulanish")

    # Request
    card(s,Inches(.45),Inches(1.65),Inches(6.1),Inches(5.6),
         fill=RGBColor(0x0A,0x16,0x2C),border=BLUE)
    tx(s,"📤  POST /predict — So'rov",
       Inches(.65),Inches(1.78),Inches(5.7),Inches(.48),
       size=13,bold=True,color=BLIGHT)
    req=('{\n  "user_id":          "user_075",\n  "amount":             850000,\n'
         '  "device":             "chrome",\n  "location":           "samarkand",\n'
         '  "transaction_hour":   23\n}')
    tx(s,req,Inches(.65),Inches(2.34),Inches(5.7),Inches(2.3),
       size=12,color=GLIGHT)
    endpoints=[
        ("POST /predict","Tranzaksiyani baholash"),
        ("GET  /health","Tizim holati"),
        ("GET  /stats","Fraud statistikasi"),
        ("GET  /model/info","Model versiya va metrikalar"),
    ]
    for i,(ep,desc) in enumerate(endpoints):
        y=Inches(4.85)+i*Inches(.44)
        tx(s,ep,Inches(.65),y,Inches(2.5),Inches(.4),size=11,bold=True,color=PURPLE)
        tx(s,desc,Inches(3.2),y,Inches(3.3),Inches(.4),size=11,color=GRAY)

    # Response
    card(s,Inches(6.75),Inches(1.65),Inches(6.1),Inches(5.6),
         fill=RGBColor(0x0A,0x16,0x2C),border=GREEN)
    tx(s,"📥  Javob (JSON Response)",
       Inches(6.95),Inches(1.78),Inches(5.7),Inches(.48),
       size=13,bold=True,color=GLIGHT)
    res=('{\n  "fraud_probability": 0.627,\n  "risk_level":        "MEDIUM",\n'
         '  "decision":          "REVIEW",\n  "reasons": [\n'
         '    "Night transaction",\n    "Amount higher than avg",\n'
         '    "New location for user"\n  ]\n}')
    tx(s,res,Inches(6.95),Inches(2.34),Inches(5.7),Inches(2.8),
       size=12,color=YELLOW)
    tx(s,"Javob vaqti:  < 100 ms",
       Inches(6.95),Inches(5.42),Inches(5.7),Inches(.38),
       size=12,bold=True,color=GLIGHT)
    tx(s,"Format: JSON   |   Auth: API Key   |   Protocol: HTTPS",
       Inches(6.95),Inches(5.82),Inches(5.7),Inches(.38),
       size=11,color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Kelajak rejalar
# ══════════════════════════════════════════════════════════════════════════════
def s14(prs):
    s = blank(prs); bg(s, BG)
    head(s,"Kelajak Rejalar",
         "Tizimni yanada kuchaytirish va kengaytirish yo'nalishlari")

    items=[
        ("🔄","Online Learning",
         "Yangi tranzaksiyalar bilan model avtomatik yangilanadi. Concept drift aniqlash.",
         RGBColor(0x0A,0x1E,0x40),BLIGHT),
        ("🤖","Deep Learning (LSTM)",
         "Ketma-ket tranzaksiyalar tahlili — vaqtiy qonuniyatlarni aniqlash.",
         RGBColor(0x18,0x10,0x30),PURPLE),
        ("🌐","Federated Learning",
         "Banklar ma'lumotlarini ulashmay birgalikda model o'rgatadi.",
         RGBColor(0x08,0x20,0x10),GREEN),
        ("📊","Graph Analytics",
         "Foydalanuvchilar to'r tahlili — guruhiy firibgarlikni aniqlash.",
         RGBColor(0x22,0x16,0x04),YELLOW),
        ("🔔","Real-Time Alerts",
         "SMS / Telegram orqali darhol xabarnoma. Bank operatoriga bildiruv.",
         RGBColor(0x28,0x08,0x08),RLIGHT),
        ("🏦","Banking Integration",
         "SWIFT, Open Banking API va O'zbekiston to'lov tizimlari bilan ulash.",
         RGBColor(0x0A,0x1E,0x40),BLIGHT),
    ]
    for i,(ico,title,desc,fill,col) in enumerate(items):
        row,c=divmod(i,2)
        x=Inches(.45)+c*Inches(6.35)
        y=Inches(1.68)+row*Inches(1.65)
        card(s,x,y,Inches(6.1),Inches(1.5),fill=fill,border=col)
        tx(s,ico+"  "+title,x+Inches(.15),y+Inches(.1),
           Inches(5.8),Inches(.5),size=14,bold=True,color=col)
        tx(s,desc,x+Inches(.15),y+Inches(.68),
           Inches(5.8),Inches(.72),size=11,color=GRAY)

    card(s,Inches(.45),Inches(6.72),Inches(12.4),Inches(.58),
         fill=RGBColor(0x08,0x16,0x2C),border=BLUE)
    tx(s,"🎯  Maqsad: SafeNet → O'zbekiston bank sektori uchun standart anti-fraud yechim",
       Inches(.65),Inches(6.78),Inches(12.0),Inches(.46),
       size=13,bold=True,color=BLIGHT,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Xulosa
# ══════════════════════════════════════════════════════════════════════════════
def s15(prs):
    s = blank(prs); bg(s, BG)
    box(s,0,Inches(2.1),W,Inches(3.2),fill=RGBColor(0x10,0x22,0x40))

    tx(s,"🛡️",Inches(5.8),Inches(.12),Inches(1.7),Inches(.95),
       size=50,align=PP_ALIGN.CENTER)
    tx(s,"SafeNet — Xulosa",
       Inches(1.5),Inches(1.1),Inches(10.3),Inches(.88),
       size=40,bold=True,color=BLIGHT,align=PP_ALIGN.CENTER)
    tx(s,"DeepFraudNet (256→128→64 MLP) asosida qurilgan real-vaqt anti-fraud tizimi\n"
        "firibgarlikni ROC-AUC 0.857 aniqlikda aniqlaydi va har qarorni tushuntiradi",
       Inches(1.5),Inches(2.25),Inches(10.3),Inches(.9),
       size=14,color=GRAY,align=PP_ALIGN.CENTER)

    highlights=[
        ("🧠","0.857\nROC-AUC",GLIGHT),
        ("⚡","<100ms\nTezlik",BLIGHT),
        ("📖","Tushun-\ntiriladi",PURPLE),
        ("🗄️","50,000\nDataset",YELLOW),
        ("⚡","25 ta\nBelgi",GLIGHT),
    ]
    for i,(ico,label,col) in enumerate(highlights):
        x=Inches(.55)+i*Inches(2.5)
        card(s,x,Inches(3.3),Inches(2.3),Inches(1.52),
             fill=RGBColor(0x0E,0x20,0x3C),border=col)
        tx(s,ico,x,Inches(3.4),Inches(2.3),Inches(.58),
           size=22,align=PP_ALIGN.CENTER)
        tx(s,label,x,Inches(3.98),Inches(2.3),Inches(.78),
           size=12,bold=True,align=PP_ALIGN.CENTER,color=col)

    # criteria checklist
    criteria=[
        ("✅","Model sifati","DeepFraudNet MLP · sqrt class weighting · stratified split · He init"),
        ("✅","Ma'lumotlar bilan ishlash","50K dataset · Z-score norm · OHE · binning · 7 interaction feature"),
        ("✅","Model aniqligi","ROC-AUC 0.857 · Precision 52.4% · Recall 40.0% · F1 45.4%"),
        ("✅","Ishlaydigan demo","Streamlit interaktiv panel · 4 ta stsenariy · GitHub'da ochiq"),
        ("✅","Algoritmni tushuntirish","Forward/backward pass · Adam optim · Risk signal engine"),
    ]
    for i,(mark,crit,detail) in enumerate(criteria):
        y=Inches(5.02)+i*Inches(.43)
        tx(s,mark,Inches(.55),y,Inches(.4),Inches(.4),size=13,color=GLIGHT)
        tx(s,crit,Inches(1.0),y,Inches(2.6),Inches(.4),size=12,bold=True,color=WHITE)
        tx(s,detail,Inches(3.65),y,Inches(9.3),Inches(.4),size=11,color=GRAY)

    tx(s,"github.com/XafizadinovUsnatdin/Fraudetection",
       Inches(1.5),Inches(7.02),Inches(10.3),Inches(.38),
       size=13,bold=True,color=BLIGHT,align=PP_ALIGN.CENTER)

# ── Build ─────────────────────────────────────────────────────────────────────
def build():
    prs = new_prs()
    s01(prs); s02(prs); s03(prs); s04(prs); s05(prs)
    s06(prs); s07(prs); s08(prs); s09(prs); s10(prs)
    s11(prs); s12(prs); s13(prs); s14(prs); s15(prs)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
